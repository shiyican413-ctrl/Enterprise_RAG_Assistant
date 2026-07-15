from pathlib import Path

from backend.ai_service.knowledge.loaders.document_loader import load_document_text


def test_load_docx_extracts_headings_paragraphs_and_tables(tmp_path: Path) -> None:
    from docx import Document

    path = tmp_path / "policy.docx"
    document = Document()
    document.add_heading("员工制度", level=1)
    document.add_paragraph("报销申请审批通过后，财务会在三个工作日内打款。")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "项目"
    table.rows[0].cells[1].text = "标准"
    document.save(path)

    text = load_document_text(path)

    assert "# 员工制度" in text
    assert "三个工作日" in text
    assert "| 项目 | 标准 |" in text


def test_load_xlsx_extracts_sheet_rows(tmp_path: Path) -> None:
    from openpyxl import Workbook

    path = tmp_path / "rules.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "差旅标准"
    sheet.append(["城市", "住宿"])
    sheet.append(["上海", "600"])
    workbook.save(path)

    text = load_document_text(path)

    assert "# 差旅标准" in text
    assert "上海 | 600" in text


def test_load_pptx_extracts_slide_text(tmp_path: Path) -> None:
    from pptx import Presentation

    path = tmp_path / "intro.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "产品介绍"
    slide.placeholders[1].text = "企业知识库支持权限隔离。"
    presentation.save(path)

    text = load_document_text(path)

    assert "[[page:1]]" in text
    assert "产品介绍" in text
    assert "权限隔离" in text
