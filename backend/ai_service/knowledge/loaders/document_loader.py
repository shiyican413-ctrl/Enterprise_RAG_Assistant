import csv
import json
from pathlib import Path


class UnsupportedDocumentType(ValueError):
    pass


def load_document_text(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")

    if suffix == ".json":
        payload = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        return json.dumps(payload, ensure_ascii=False, indent=2)

    if suffix == ".csv":
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as file:
            rows = list(csv.reader(file))
        return _markdown_table(rows)

    if suffix == ".pdf":
        return _load_pdf(path)

    if suffix == ".docx":
        return _load_docx(path)

    if suffix == ".xlsx":
        return _load_xlsx(path)

    if suffix == ".pptx":
        return _load_pptx(path)

    raise UnsupportedDocumentType(f"Unsupported document type: {suffix}")


def _load_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise UnsupportedDocumentType(
            "PDF support requires installing optional dependency: pypdf"
        ) from exc

    reader = PdfReader(str(path))
    pages: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        pages.append(f"[[page:{index}]]\n{page_text}")
    return "\n\n".join(pages).strip()


def _load_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise UnsupportedDocumentType(
            "DOCX support requires installing dependency: python-docx"
        ) from exc

    document = Document(str(path))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style.startswith("heading"):
            level = "".join(char for char in style if char.isdigit()) or "2"
            parts.append(f"{'#' * min(int(level), 6)} {text}")
        else:
            parts.append(text)

    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts).strip()


def _load_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise UnsupportedDocumentType(
            "XLSX support requires installing dependency: openpyxl"
        ) from exc

    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    sheets: list[str] = []
    for sheet in workbook.worksheets:
        raw_rows: list[list[object]] = []
        for row in sheet.iter_rows(values_only=True):
            raw_rows.append(list(row))
        table = _markdown_table(raw_rows)
        sheets.append("\n\n".join(part for part in [f"# {sheet.title}", table] if part))
    workbook.close()
    return "\n\n".join(sheets).strip()


def _load_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise UnsupportedDocumentType(
            "PPTX support requires installing dependency: python-pptx"
        ) from exc

    presentation = Presentation(str(path))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = [f"[[page:{index}]]", f"# Slide {index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text or "").strip()
                if text:
                    texts.append(text)
        slides.append("\n".join(texts))
    return "\n\n".join(slides).strip()


def _markdown_table(rows: list[list[object]]) -> str:
    cleaned = [
        ["" if cell is None else str(cell).strip().replace("\n", " ") for cell in row]
        for row in rows
        if any("" if cell is None else str(cell).strip() for cell in row)
    ]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    normalized = [row + [""] * (width - len(row)) for row in cleaned]
    lines = ["| " + " | ".join(row) + " |" for row in normalized]
    if len(lines) > 1:
        separator = "| " + " | ".join("---" for _ in range(width)) + " |"
        lines.insert(1, separator)
    return "\n".join(lines)
